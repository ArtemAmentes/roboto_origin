# -*- coding: utf-8 -*-
"""
Презентация: Национальная платформа робототехники.
Формат государственного проекта для Москвы.
Стиль: холодный хай-тек металлик.
Запуск: python create_presentation.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DIR = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(DIR, "presentation_roboty_subsidy.pptx")
IMG_ROBOT = os.path.join(DIR, "photo_2026-02-19 11.20.26.jpeg")
IMG_SLIDE2 = os.path.join(DIR, "robots_components_transfer.png")

SW = Inches(13.333)
SH = Inches(7.5)

BG        = RGBColor(0xE8, 0xEC, 0xF1)
CARD_CLR  = RGBColor(0xF5, 0xF7, 0xFA)
EDGE      = RGBColor(0xC4, 0xCC, 0xD8)
DARK      = RGBColor(0x1A, 0x1F, 0x2E)
MID       = RGBColor(0x4A, 0x55, 0x68)
BLUE      = RGBColor(0x1B, 0x6E, 0xC2)
DBLUE     = RGBColor(0x0F, 0x3D, 0x6B)
GREEN     = RGBColor(0x0E, 0x7C, 0x5A)
ARROW_CLR = RGBColor(0x8A, 0xA8, 0xCC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SCHEME_BG = RGBColor(0xDB, 0xE4, 0xF0)

F = "Arial"
MX = Inches(0.7)
GAP = Inches(0.25)
PAD = Inches(0.3)


# --- helpers ---

def set_bg(s):
    f = s.background.fill; f.solid(); f.fore_color.rgb = BG

def top_bar(s):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.05))
    b.fill.solid(); b.fill.fore_color.rgb = BLUE; b.line.fill.background()

def card(s, l, t, w, h, fill=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill or CARD_CLR
    sh.line.color.rgb = EDGE; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

def accent_card(s, l, t, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = DBLUE
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh

def scheme_block(s, l, t, w, h, label, clr=BLUE, bg_clr=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg_clr or SCHEME_BG
    sh.line.color.rgb = BLUE; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    txt(s, label, l + Inches(0.1), t + Inches(0.05), w - Inches(0.2), h - Inches(0.1),
        sz=Pt(10), clr=clr, bold=True, al=PP_ALIGN.CENTER)

def arrow_r(s, l, t, w=Inches(0.4), h=Inches(0.28)):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = ARROW_CLR
    sh.line.fill.background(); sh.shadow.inherit = False

def arrow_d(s, l, t, w=Inches(0.28), h=Inches(0.35)):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = ARROW_CLR
    sh.line.fill.background(); sh.shadow.inherit = False

def txt(s, text, l, t, w, h, sz=Pt(14), clr=DARK, bold=False, al=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = sz; p.font.name = F
    p.font.bold = bold; p.font.color.rgb = clr; p.alignment = al
    return bx

def lines(s, items, l, t, w, h, sz=Pt(13), clr=DARK, sp=Pt(14)):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, ln in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.font.size = sz; p.font.name = F
        p.font.color.rgb = clr; p.space_after = sp
    return bx

def fit_pic(slide, path, left, top, max_w, max_h):
    """Вставка картинки с сохранением пропорций (без растягивания)."""
    pic = slide.shapes.add_picture(path, left, top, width=max_w)
    if pic.height > max_h:
        el = pic._element; el.getparent().remove(el)
        pic = slide.shapes.add_picture(path, left, top, height=max_h)
        if pic.width > max_w:
            el = pic._element; el.getparent().remove(el)
            pic = slide.shapes.add_picture(path, left, top, width=max_w, height=max_h)
    pic.left = left + (max_w - pic.width) // 2
    pic.top = top + (max_h - pic.height) // 2
    return pic


# =====================================================================
#  1. ТИТУЛЬНЫЙ СЛАЙД
# =====================================================================
def slide_01(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "amentes  \u00B7  Социальный код", MX, Inches(0.4), Inches(6), Inches(0.3),
        sz=Pt(12), clr=MID)

    txt(s, "Национальная платформа\nроссийской робототехники",
        MX, Inches(1.0), Inches(6.5), Inches(1.8),
        sz=Pt(42), bold=True, clr=DARK)

    txt(s, "Разработка отечественного робота-андроида\nс открытым SDK, модульной компонентной базой\n"
           "и программным обеспечением для автоматизации\nпроизводств, социальных учреждений и быта",
        MX, Inches(3.1), Inches(6), Inches(1.5), sz=Pt(16), clr=MID)

    txt(s, "Проект для Правительства Москвы",
        MX, Inches(5.0), Inches(5), Inches(0.4), sz=Pt(14), clr=BLUE, bold=True)

    if os.path.isfile(IMG_ROBOT):
        fit_pic(s, IMG_ROBOT, Inches(7.8), Inches(0.5), Inches(5.2), Inches(6.5))


# =====================================================================
#  2. ВЫЗОВ
# =====================================================================
def slide_02(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Вызов: зависимость от импорта и дефицит кадров",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    problems = [
        ("Импортозависимость",
         "Робототехника в России на 95% состоит из зарубежных компонентов и ПО. "
         "Санкции ограничивают доступ к ключевым технологиям. "
         "Отсутствует единая платформа для перепрошивки и адаптации зарубежных моделей."),
        ("Дефицит рабочих рук",
         "Соцучреждения, ЖКХ, производства испытывают хронический кадровый голод. "
         "Рутинные задачи, для которых раньше нужен был человек, могут выполнять роботы-андроиды. "
         "Помощь по дому, обслуживание в учреждениях, работа на конвейере."),
        ("Отсутствие экосистемы",
         "Нет открытого российского SDK для управления робототехникой. "
         "Нет единой номенклатуры узлов и агрегатов. "
         "Нет системной подготовки кадров для отрасли."),
    ]
    cw = (SW - 2*MX - 2*GAP) / 3
    ch = Inches(4.5); top = Inches(1.2)
    for i, (title, body) in enumerate(problems):
        left = MX + i * (cw + GAP)
        card(s, left, top, cw, ch)
        txt(s, title, left + PAD, top + Inches(0.2), cw - 2*PAD, Inches(0.4),
            sz=Pt(16), clr=BLUE, bold=True)
        txt(s, body, left + PAD, top + Inches(0.75), cw - 2*PAD, Inches(3.5),
            sz=Pt(13), clr=DARK)

    txt(s, "Решение — создание национальной платформы робототехники с открытым ПО, модульным железом и системой подготовки кадров",
        MX, Inches(6.1), SW - 2*MX, Inches(0.5), sz=Pt(14), clr=MID)


# =====================================================================
#  3. ТРИ СТОЛПА
# =====================================================================
def slide_03(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Решение: три столпа национальной платформы",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    pillars = [
        ("01  Железо",
         "Стратегия постепенного перехода на российскую компонентную базу "
         "с разработкой универсальной модульной системы замены частей и агрегатов. "
         "Импорт мировых образцов андроидной техники из Китая, Индии, США, Европы, Японии — "
         "изучение структуры, разработка российских аналогов, создание ЗИП."),
        ("02  Софт",
         "Открытый программный код для прошивки моделей робота под домашние задачи, "
         "автоматизацию производств и социальные учреждения. "
         "LLM-агент для быстрого переноса человеческого языка тех. задания в алгоритм движения робота. "
         "Централизованное обновление, а также закрытый контур управления."),
        ("03  Образование",
         "Обучение / колледж для молодёжи по робототехнике и сборке конкретных моделей "
         "роботов будущего российского производства. "
         "Создание сборочных мастерских для сборки, ремонта и обслуживания робототехники. "
         "Подготовка кадров для промышленной эксплуатации и сервиса роботов."),
    ]
    cw = (SW - 2*MX - 2*GAP) / 3
    ch = Inches(4.2); top = Inches(1.2)
    for i, (title, body) in enumerate(pillars):
        left = MX + i * (cw + GAP)
        card(s, left, top, cw, ch)
        txt(s, title, left + PAD, top + Inches(0.2), cw - 2*PAD, Inches(0.4),
            sz=Pt(18), clr=BLUE, bold=True)
        txt(s, body, left + PAD, top + Inches(0.8), cw - 2*PAD, Inches(3.0),
            sz=Pt(13), clr=DARK)

    # --- Схема связей: три столпа соединяются стрелками в единую платформу ---
    row_y = Inches(5.7)
    bw = Inches(1.8); bh = Inches(0.45)
    aw = Inches(0.35); ah = Inches(0.25)
    centers = [MX + cw/2, MX + cw + GAP + cw/2, MX + 2*(cw + GAP) + cw/2]

    for i, label in enumerate(["Железо", "Софт", "Образование"]):
        scheme_block(s, centers[i] - bw/2, row_y, bw, bh, label)

    for i in range(2):
        arrow_r(s, centers[i] + bw/2 + Inches(0.05), row_y + Inches(0.08),
                w=centers[i+1] - bw/2 - (centers[i] + bw/2) - Inches(0.1), h=ah)

    txt(s, "Открытый SDK  \u00B7  Открытая номенклатура  \u00B7  Открытое образование",
        MX, Inches(6.4), SW - 2*MX, Inches(0.4),
        sz=Pt(15), clr=BLUE, bold=True, al=PP_ALIGN.CENTER)


# =====================================================================
#  4. СОФТ: SDK И АВТОМАТИЗАЦИЯ
# =====================================================================
def slide_04(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Программное обеспечение: SDK и автоматизация",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    col_w = (SW - 2*MX - GAP) / 2
    card_h = Inches(4.2)

    card(s, MX, Inches(1.1), col_w, card_h)
    txt(s, "Российский SDK для робототехники",
        MX + PAD, Inches(1.3), col_w - 2*PAD, Inches(0.4), sz=Pt(16), clr=BLUE, bold=True)
    lines(s, [
        "Открытый программный код для прошивки моделей робота под любые задачи: домашние, промышленные, социальные",
        "LLM-агент для быстрого переноса человеческого языка технического задания в алгоритм движения робота — оператору не нужно быть программистом",
        "Централизованное обновление прошивок через защищённый канал, а также закрытый контур для критической инфраструктуры",
        "Единый API для интеграции роботов в существующие системы управления производством (MES, SCADA, ERP)",
        "Модульная архитектура: подключение новых датчиков, актуаторов и периферии без перекомпиляции ядра",
    ], MX + PAD, Inches(1.85), col_w - 2*PAD, Inches(3.2), sz=Pt(12), clr=DARK, sp=Pt(12))

    r = MX + col_w + GAP
    card(s, r, Inches(1.1), col_w, card_h)
    txt(s, "Полная автоматизация производств",
        r + PAD, Inches(1.3), col_w - 2*PAD, Inches(0.4), sz=Pt(16), clr=BLUE, bold=True)
    lines(s, [
        "Софт для программирования роботов на производственных линиях: сборка, сортировка, упаковка, контроль качества",
        "Интеграция с конвейерными системами и промышленными контроллерами (PLC)",
        "Визуальный конструктор производственных сценариев — настройка без кода, на основе LLM-агента и голосовых команд",
        "Автоматическая адаптация алгоритмов к смене номенклатуры изделий и переналадке линий",
        "Мониторинг состояния роботов в реальном времени: предиктивное обслуживание, диагностика, отчёты для руководства предприятия",
        "Применение в производствах г. Москвы: машиностроение, пищевая промышленность, логистика, фармацевтика",
    ], r + PAD, Inches(1.85), col_w - 2*PAD, Inches(3.2), sz=Pt(12), clr=DARK, sp=Pt(10))

    # --- Схема: как работает SDK (горизонтальный поток) ---
    flow_y = Inches(5.6)
    steps = ["ТЗ на языке\nчеловека", "LLM-агент", "Алгоритм\nдвижения", "Робот\nна линии", "Готовая\nпродукция"]
    bw = Inches(1.6); bh = Inches(0.65)
    total = len(steps) * bw + (len(steps) - 1) * Inches(0.55)
    start_x = (SW - total) / 2

    for i, label in enumerate(steps):
        x = start_x + i * (bw + Inches(0.55))
        bg_c = DBLUE if i == 1 else SCHEME_BG
        clr_c = WHITE if i == 1 else BLUE
        scheme_block(s, x, flow_y, bw, bh, label, clr=clr_c, bg_clr=bg_c)
        if i < len(steps) - 1:
            arrow_r(s, x + bw + Inches(0.05), flow_y + Inches(0.18),
                    w=Inches(0.45), h=Inches(0.25))

    txt(s, "Схема: от задачи на человеческом языке до результата на производстве",
        MX, Inches(6.5), SW - 2*MX, Inches(0.4), sz=Pt(11), clr=MID, al=PP_ALIGN.CENTER)


# =====================================================================
#  5. ЖЕЛЕЗО: МОДУЛЬНАЯ ПЛАТФОРМА
# =====================================================================
def slide_05(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Аппаратная платформа: модульность и импортозамещение",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    stages = [
        ("Импорт образцов",
         "Закупка мировых образцов андроидной техники: "
         "Китай, Индия, США, Европа, Япония. "
         "Разборка, анализ структуры, каталогизация компонентов."),
        ("Разработка SDK",
         "Разработка российского SDK программного обеспечения "
         "для перепрошивки и управления роботами. "
         "Универсальный интерфейс для всех моделей."),
        ("Замена компонентов",
         "Постепенная замена компонентов, узлов и агрегатов "
         "на российские аналоги. Универсальная модульная система замены частей."),
        ("Сборочные мастерские",
         "Создание мастерских для сборки, ремонта и обслуживания "
         "робототехники. Формирование запасов ЗИП для всех моделей."),
    ]
    cw = (SW - 2*MX - 3*GAP) / 4
    ch = Inches(3.0); top = Inches(1.15)

    for i, (title, body) in enumerate(stages):
        left = MX + i * (cw + GAP)
        card(s, left, top, cw, ch)
        txt(s, f"0{i+1}", left + PAD, top + Inches(0.15), Inches(0.5), Inches(0.35),
            sz=Pt(20), clr=BLUE, bold=True)
        txt(s, title, left + PAD, top + Inches(0.55), cw - 2*PAD, Inches(0.35),
            sz=Pt(14), clr=BLUE, bold=True)
        txt(s, body, left + PAD, top + Inches(0.95), cw - 2*PAD, Inches(1.8),
            sz=Pt(12), clr=DARK)
        if i < 3:
            arrow_r(s, left + cw + Inches(0.01), top + Inches(1.2))

    if os.path.isfile(IMG_SLIDE2):
        fit_pic(s, IMG_SLIDE2, Inches(2.0), Inches(4.4), Inches(9.3), Inches(2.8))


# =====================================================================
#  6. ПРИМЕНЕНИЕ: ТРИ СФЕРЫ
# =====================================================================
def slide_06(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Сферы применения платформы",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    # --- Схема: SDK как центральный хаб → три сферы ---
    hub_w = Inches(2.5); hub_h = Inches(0.5)
    hub_x = (SW - hub_w) / 2; hub_y = Inches(1.1)
    scheme_block(s, hub_x, hub_y, hub_w, hub_h, "Единый SDK — платформа", clr=WHITE, bg_clr=DBLUE)

    cw = (SW - 2*MX - 2*GAP) / 3
    sphere_labels = ["Производство", "Госучреждения", "Быт"]
    for i, lbl in enumerate(sphere_labels):
        center_x = MX + i * (cw + GAP) + cw / 2
        arrow_d(s, center_x - Inches(0.14), hub_y + hub_h + Inches(0.02),
                w=Inches(0.28), h=Inches(0.3))

    areas = [
        ("Производство и промышленность",
         [
             "Полная автоматизация производственных линий: сборка, сортировка, упаковка, контроль качества",
             "Интеграция с существующими конвейерными системами и промышленными контроллерами",
             "Визуальный конструктор сценариев без программирования — через LLM-агент",
             "Предиктивное обслуживание оборудования и мониторинг состояния роботов",
             "Применение: машиностроение, пищевая промышленность, логистика, фармацевтика",
         ]),
        ("Социальные учреждения",
         [
             "Помощь персоналу в больницах, поликлиниках, домах престарелых, МФЦ",
             "Навигация посетителей, выдача документов, первичный приём обращений",
             "Персонализация через CVCODE — робот «знает» гражданина и адаптирует общение",
             "Распознавание эмоций через ANIEMORE — эмпатичное взаимодействие на русском языке",
         ]),
        ("Бытовое применение",
         [
             "Помощь по дому: уборка, контроль, рутинные задачи",
             "Управление голосом на русском языке через LLM-агент",
             "Открытый код прошивки — настройка робота под конкретные домашние задачи",
             "Снижение нагрузки на семьи с пожилыми и маломобильными членами",
         ]),
    ]
    ch = Inches(4.5); card_top = Inches(2.0)

    for i, (title, bullets) in enumerate(areas):
        left = MX + i * (cw + GAP)
        card(s, left, card_top, cw, ch)
        txt(s, title, left + PAD, card_top + Inches(0.15), cw - 2*PAD, Inches(0.45),
            sz=Pt(15), clr=BLUE, bold=True)
        lines(s, bullets, left + PAD, card_top + Inches(0.65), cw - 2*PAD, Inches(3.5),
              sz=Pt(12), clr=DARK, sp=Pt(11))

    txt(s, "Единый SDK — один софт для всех сфер применения",
        MX, Inches(6.8), SW - 2*MX, Inches(0.4), sz=Pt(14), clr=MID, al=PP_ALIGN.CENTER)


# =====================================================================
#  7. CVCODE + ANIEMORE
# =====================================================================
def slide_07(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Технологии персонализации и эмпатии",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    cw = (SW - 2*MX - GAP) / 2
    ch = Inches(3.6); top = Inches(1.1)

    card(s, MX, top, cw, ch)
    txt(s, "CVCODE", MX + PAD, top + Inches(0.2), cw - 2*PAD, Inches(0.4),
        sz=Pt(20), clr=BLUE, bold=True)
    lines(s, [
        "Источник персонализации и гражданской идентичности робота и его ИИ",
        "Команда разработала CVCODE, который позволит сделать роботов более «человечными» для работы в социальных учреждениях",
        "Робот распознаёт конкретного гражданина, помнит контекст взаимодействия, адаптирует поведение",
        "Интеграция с государственными системами идентификации и учёта обращений граждан",
    ], MX + PAD, top + Inches(0.7), cw - 2*PAD, Inches(2.7), sz=Pt(13), clr=DARK, sp=Pt(12))

    r = MX + cw + GAP
    card(s, r, top, cw, ch)
    txt(s, "ANIEMORE", r + PAD, top + Inches(0.2), cw - 2*PAD, Inches(0.4),
        sz=Pt(20), clr=BLUE, bold=True)
    lines(s, [
        "Система распознавания эмоций в голосе — лучшая в русском языке",
        "Разработана для повышения эмпатии роботов при взаимодействии с людьми",
        "Робот определяет эмоциональное состояние человека и корректирует тон, скорость речи и формулировки",
        "Особенно важно для работы в медицинских и социальных учреждениях, взаимодействие с пожилыми людьми",
    ], r + PAD, top + Inches(0.7), cw - 2*PAD, Inches(2.7), sz=Pt(13), clr=DARK, sp=Pt(12))

    # --- Схема взаимодействия: Гражданин → CVCODE → ANIEMORE → Робот ---
    flow_y = Inches(5.05)
    steps = [
        ("Гражданин", SCHEME_BG, BLUE),
        ("CVCODE\nидентификация", DBLUE, WHITE),
        ("ANIEMORE\nэмоции", DBLUE, WHITE),
        ("Робот: адаптивное\nповедение", SCHEME_BG, BLUE),
    ]
    bw = Inches(2.2); bh = Inches(0.7)
    total = len(steps) * bw + (len(steps) - 1) * Inches(0.55)
    start_x = (SW - total) / 2
    for i, (label, bg_c, fg_c) in enumerate(steps):
        x = start_x + i * (bw + Inches(0.55))
        scheme_block(s, x, flow_y, bw, bh, label, clr=fg_c, bg_clr=bg_c)
        if i < len(steps) - 1:
            arrow_r(s, x + bw + Inches(0.06), flow_y + Inches(0.2),
                    w=Inches(0.43), h=Inches(0.25))

    txt(s, "Схема: от идентификации гражданина до адаптивного поведения робота",
        MX, Inches(5.95), SW - 2*MX, Inches(0.35), sz=Pt(11), clr=MID, al=PP_ALIGN.CENTER)

    txt(s, "Гражданская идентичность + эмпатия = доверие граждан к роботизированным сервисам в государственных учреждениях",
        MX, Inches(6.4), SW - 2*MX, Inches(0.5), sz=Pt(14), clr=MID, al=PP_ALIGN.CENTER)


# =====================================================================
#  8. СХЕМА ИНТЕГРАЦИИ ЭКОСИСТЕМЫ
# =====================================================================
def slide_08(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Схема интеграции: единая экосистема",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    top_w = SW - 2*MX
    accent_card(s, MX, Inches(1.1), top_w, Inches(0.7))
    txt(s, "Открытая платформа: SDK  \u00B7  Номенклатура узлов  \u00B7  Образование  \u00B7  Опенсорс-сообщество",
        MX + PAD, Inches(1.2), top_w - 2*PAD, Inches(0.5),
        sz=Pt(15), clr=WHITE, bold=True, al=PP_ALIGN.CENTER)

    mod_w = (SW - 2*MX - 2*GAP) / 3
    centers = [MX + mod_w/2, MX + mod_w + GAP + mod_w/2, MX + 2*(mod_w + GAP) + mod_w/2]

    # Стрелки: платформа → модули
    for cx in centers:
        arrow_d(s, cx - Inches(0.14), Inches(1.85))

    # Модули
    modules = [
        ("SDK робототехники",
         "Прошивка  \u00B7  LLM-агент\nУправление движением\nИнтеграция MES/SCADA"),
        ("CVCODE",
         "Персонализация\nГражданская идентичность\nАдаптивное поведение"),
        ("ANIEMORE",
         "Распознавание эмоций\nЭмпатия на русском языке\nКоррекция поведения"),
    ]
    mod_h = Inches(1.7); mod_top = Inches(2.3)
    for i, (title, body) in enumerate(modules):
        left = MX + i * (mod_w + GAP)
        card(s, left, mod_top, mod_w, mod_h)
        txt(s, title, left + PAD, mod_top + Inches(0.1), mod_w - 2*PAD, Inches(0.3),
            sz=Pt(14), clr=BLUE, bold=True)
        txt(s, body, left + PAD, mod_top + Inches(0.45), mod_w - 2*PAD, Inches(1.1),
            sz=Pt(12), clr=DARK)

    # Горизонтальные стрелки между модулями (взаимосвязь)
    for i in range(2):
        l_edge = MX + (i+1) * mod_w + i * GAP + Inches(0.02)
        arrow_r(s, l_edge, mod_top + Inches(0.7), w=GAP - Inches(0.04), h=Inches(0.22))

    # Стрелки: модули → применение
    for cx in centers:
        arrow_d(s, cx - Inches(0.14), Inches(4.05))

    # Применение
    apps = [
        ("Производство",
         "Автоматизация линий\nМашиностроение, пищепром\nЛогистика, фармацевтика"),
        ("Госучреждения",
         "Больницы, МФЦ\nНавигация, приём обращений\nДеп. здравоохранения Москвы"),
        ("Быт",
         "Помощь по дому\nУход за пожилыми\nУправление голосом"),
    ]
    app_top = Inches(4.5)
    for i, (title, body) in enumerate(apps):
        left = MX + i * (mod_w + GAP)
        card(s, left, app_top, mod_w, Inches(1.6))
        txt(s, title, left + PAD, app_top + Inches(0.1), mod_w - 2*PAD, Inches(0.3),
            sz=Pt(14), clr=BLUE, bold=True)
        txt(s, body, left + PAD, app_top + Inches(0.45), mod_w - 2*PAD, Inches(1.0),
            sz=Pt(12), clr=DARK)

    txt(s, "Стратегия открытой лаборатории и опенсорс для проектов роботов России. "
           "Опора на прозрачность, сообщество и переиспользование решений.",
        MX, Inches(6.5), SW - 2*MX, Inches(0.5), sz=Pt(12), clr=MID)


# =====================================================================
#  9. ФИНАНСИРОВАНИЕ — ЭТАП 1: МОСКВА
# =====================================================================
def slide_09(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Этап 1: экспериментальная лаборатория в Москве",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    # Шапка
    accent_card(s, MX, Inches(1.05), SW - 2*MX, Inches(0.5))
    txt(s, "Бюджет: 1 млрд \u20BD  \u00B7  Создание в Москве экспериментальной площадки мирового уровня",
        MX + Inches(0.2), Inches(1.1), SW - 2*MX - Inches(0.4), Inches(0.4),
        sz=Pt(15), clr=WHITE, bold=True, al=PP_ALIGN.CENTER)

    col_w = (SW - 2*MX - GAP) / 2

    # Левая карточка: что делаем
    card(s, MX, Inches(1.75), col_w, Inches(4.8))
    txt(s, "Что получает Москва за 1 млрд \u20BD",
        MX + PAD, Inches(1.9), col_w - 2*PAD, Inches(0.4),
        sz=Pt(16), clr=BLUE, bold=True)
    lines(s, [
        "Импорт лучших мировых образцов андроидной техники из Китая, Индии, США, Европы и Японии — полная разборка, реверс-инжиниринг, каталогизация компонентов",
        "Разработка российского SDK программного обеспечения для перепрошивки и управления импортными роботами — единая платформа для всех моделей",
        "Экспериментальная сборка роботов в Москве на базе иностранных компонентов с постепенной заменой на отечественные аналоги",
        "Формирование проектной команды инженеров и разработчиков под проект в Москве (кадры МФТИ, Бауман, ФИЦ ИУ РАН)",
        "Создание первой в России сборочной мастерской робототехники в Москве — ремонт, обслуживание, ЗИП",
        "Запуск пилотных внедрений роботов на московских предприятиях и в учреждениях Правительства Москвы",
    ], MX + PAD, Inches(2.45), col_w - 2*PAD, Inches(3.9), sz=Pt(12), clr=DARK, sp=Pt(10))

    # Правая карточка: результат для Москвы
    r = MX + col_w + GAP
    card(s, r, Inches(1.75), col_w, Inches(4.8))
    txt(s, "Результат Этапа 1 для Москвы",
        r + PAD, Inches(1.9), col_w - 2*PAD, Inches(0.4),
        sz=Pt(16), clr=GREEN, bold=True)
    lines(s, [
        "Москва получает действующую экспериментальную лабораторию робототехники — первую в стране",
        "Полностью рабочий прототип российского робота-андроида на базе изученных иностранных компонентов",
        "Открытый SDK, готовый к применению на производствах и в госучреждениях Москвы",
        "Собранная и обученная команда — ядро будущей отрасли робототехники в России",
        "Первые пилотные роботы в учреждениях Департамента здравоохранения Москвы и на московских предприятиях",
        "Москва становится центром компетенций в робототехнике — аналог Силиконовой долины для российских роботов",
    ], r + PAD, Inches(2.45), col_w - 2*PAD, Inches(3.9), sz=Pt(12), clr=DARK, sp=Pt(10))

    txt(s, "Этап 1 — фундамент: экспериментальный подход, реверс-инжиниринг, формирование команды, первые пилоты в Москве",
        MX, Inches(6.8), SW - 2*MX, Inches(0.4), sz=Pt(13), clr=MID, al=PP_ALIGN.CENTER)


# =====================================================================
#  10. ФИНАНСИРОВАНИЕ — ЭТАП 2: МАСШТАБИРОВАНИЕ
# =====================================================================
def slide_10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Этап 2: промышленное производство и масштабирование",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    accent_card(s, MX, Inches(1.05), SW - 2*MX, Inches(0.5))
    txt(s, "Бюджет: 9 млрд \u20BD  \u00B7  Москва — столица российской робототехники, масштабирование на всю Россию",
        MX + Inches(0.2), Inches(1.1), SW - 2*MX - Inches(0.4), Inches(0.4),
        sz=Pt(15), clr=WHITE, bold=True, al=PP_ALIGN.CENTER)

    cw = (SW - 2*MX - 2*GAP) / 3
    ch = Inches(4.6); top = Inches(1.75)

    # Колонка 1: Производство
    card(s, MX, top, cw, ch)
    txt(s, "Производство в Москве",
        MX + PAD, top + Inches(0.15), cw - 2*PAD, Inches(0.4),
        sz=Pt(15), clr=BLUE, bold=True)
    lines(s, [
        "Запуск в Москве производственных линий компонентов, узлов и агрегатов роботов — на основе реверс-инжиниринга иностранных образцов",
        "Запуск производственных линий по сборке готовых роботов-андроидов российского производства в Москве",
        "Создание в Москве полного цикла: от компонента до готового робота",
        "Формирование запасов ЗИП для обеспечения эксплуатации роботов по всей стране",
        "Москва — центр серийного производства, экспорт технологий в регионы",
    ], MX + PAD, top + Inches(0.65), cw - 2*PAD, Inches(3.7), sz=Pt(12), clr=DARK, sp=Pt(10))

    # Колонка 2: Кадры и образование
    c2 = MX + cw + GAP
    card(s, c2, top, cw, ch)
    txt(s, "Кадры и образование",
        c2 + PAD, top + Inches(0.15), cw - 2*PAD, Inches(0.4),
        sz=Pt(15), clr=BLUE, bold=True)
    lines(s, [
        "Создание на базе ведущих московских ВУЗов кафедр робототехники (МФТИ, Бауманка, МАИ и др.)",
        "Подготовка инженеров-робототехников, программистов SDK, специалистов по ИИ для роботов",
        "Обучение / колледж для молодёжи по сборке конкретных моделей роботов будущего российского производства",
        "Программы стажировок на московском заводе робототехники для студентов из регионов",
        "Поиск талантов по всей России с помощью системы CVCODE — лучшие кадры стекаются в Москву",
    ], c2 + PAD, top + Inches(0.65), cw - 2*PAD, Inches(3.7), sz=Pt(12), clr=DARK, sp=Pt(10))

    # Колонка 3: Масштабирование
    c3 = MX + 2*(cw + GAP)
    card(s, c3, top, cw, ch)
    txt(s, "Масштабирование из Москвы",
        c3 + PAD, top + Inches(0.15), cw - 2*PAD, Inches(0.4),
        sz=Pt(15), clr=BLUE, bold=True)
    lines(s, [
        "Открытие филиалов московского завода робототехники в ключевых регионах России",
        "Региональные сборочные мастерские и сервисные центры — по стандартам московской площадки",
        "Трансфер технологий и методологий из Москвы в регионы",
        "Обеспечение роботами региональных производств, соцучреждений и населения",
        "Москва — управляющий центр национальной сети робототехники, регионы — точки роста",
    ], c3 + PAD, top + Inches(0.65), cw - 2*PAD, Inches(3.7), sz=Pt(12), clr=DARK, sp=Pt(10))

    txt(s, "Москва — Силиконовая долина российской робототехники: от лаборатории до серийного производства и экспорта технологий в регионы",
        MX, Inches(6.65), SW - 2*MX, Inches(0.5), sz=Pt(14), clr=BLUE, bold=True, al=PP_ALIGN.CENTER)


# =====================================================================
#  11. ИТОГИ: ЧТО ПОЛУЧАЕТ МОСКВА
# =====================================================================
def slide_11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Итого: что получает Москва",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    # Таймлайн
    tl_y = Inches(1.05); e_w = (SW - 2*MX - Inches(0.5)) / 2
    accent_card(s, MX, tl_y, e_w, Inches(0.5))
    txt(s, "Этап 1  \u00B7  1 млрд \u20BD", MX + Inches(0.15), tl_y + Inches(0.08),
        e_w - Inches(0.3), Inches(0.35), sz=Pt(14), clr=WHITE, bold=True, al=PP_ALIGN.CENTER)
    arrow_r(s, MX + e_w + Inches(0.03), tl_y + Inches(0.1), w=Inches(0.4), h=Inches(0.25))
    accent_card(s, MX + e_w + Inches(0.5), tl_y, e_w, Inches(0.5))
    txt(s, "Этап 2  \u00B7  9 млрд \u20BD", MX + e_w + Inches(0.65), tl_y + Inches(0.08),
        e_w - Inches(0.3), Inches(0.35), sz=Pt(14), clr=WHITE, bold=True, al=PP_ALIGN.CENTER)

    # Результаты в три колонки
    cw = (SW - 2*MX - 2*GAP) / 3
    ch = Inches(3.5); top = Inches(1.8)

    results = [
        ("Открытый SDK",
         [
             "Открытый программный пакет SDK для управления ПО робототехники",
             "Единая платформа: от бытового робота до промышленного конвейера",
             "LLM-агент: управление роботами на человеческом языке",
             "Централизованное обновление и закрытый контур для критической инфраструктуры",
         ]),
        ("Открытая номенклатура",
         [
             "Полный каталог узлов и агрегатов робототехники с российскими аналогами",
             "Производственные линии компонентов в Москве",
             "Серийное производство роботов-андроидов",
             "Запас ЗИП и сервисная инфраструктура",
         ]),
        ("Открытое образование",
         [
             "Кафедры робототехники в московских ВУЗах",
             "Колледж по сборке и обслуживанию роботов",
             "Система поиска талантов (CVCODE) по всей России",
             "Москва — кузница кадров для новой отрасли",
         ]),
    ]

    for i, (title, bullets) in enumerate(results):
        left = MX + i * (cw + GAP)
        card(s, left, top, cw, ch)
        txt(s, title, left + PAD, top + Inches(0.15), cw - 2*PAD, Inches(0.4),
            sz=Pt(16), clr=BLUE, bold=True)
        lines(s, bullets, left + PAD, top + Inches(0.65), cw - 2*PAD, Inches(2.6),
              sz=Pt(12), clr=DARK, sp=Pt(12))

    # Итоговая фраза
    accent_card(s, MX, Inches(5.6), SW - 2*MX, Inches(0.7))
    txt(s, "10 млрд \u20BD — и Москва становится мировым центром робототехники: собственное производство, "
           "собственное ПО, собственные кадры. Российская Силиконовая долина роботов.",
        MX + PAD, Inches(5.65), SW - 2*MX - 2*PAD, Inches(0.6),
        sz=Pt(15), clr=WHITE, bold=True, al=PP_ALIGN.CENTER)


# =====================================================================
#  12. КОМАНДА И ОПЫТ
# =====================================================================
def slide_12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); top_bar(s)

    txt(s, "Команда и подтверждённый опыт",
        MX, Inches(0.3), SW - 2*MX, Inches(0.6), sz=Pt(30), bold=True)

    cw = (SW - 2*MX - GAP) / 2

    ch1 = Inches(2.6); t1 = Inches(1.1)
    card(s, MX, t1, cw, ch1)
    txt(s, "Социальный код", MX + PAD, t1 + Inches(0.15), cw - 2*PAD, Inches(0.35),
        sz=Pt(17), clr=BLUE, bold=True)
    lines(s, [
        "Компания Социальный код (Екатеринбург) обладает кадровым и методологическим капиталом для реализации проекта",
        "Руководитель: Аментес Артём Важаевич — выпускник МФТИ, аспирант ФИЦ ИУ РАН по направлению ИИ",
        "Команда инженеров МГТУ им. Н. Э. Баумана",
    ], MX + PAD, t1 + Inches(0.6), cw - 2*PAD, Inches(1.8), sz=Pt(13), clr=DARK, sp=Pt(12))

    card(s, MX + cw + GAP, t1, cw, ch1)
    txt(s, "Civicode (Сивикод)", MX + cw + GAP + PAD, t1 + Inches(0.15),
        cw - 2*PAD, Inches(0.35), sz=Pt(17), clr=BLUE, bold=True)
    lines(s, [
        "Всероссийское открытое ПО для оцифровки личности человека",
        "Уже внедряется в Правительстве Москвы (Департамент здравоохранения Москвы), заключены первые госконтракты",
        "Сивикод сопровождал знаковые федеральные проекты: «Школьная проектная олимпиада», «Стройотряды России» и др.",
    ], MX + cw + GAP + PAD, t1 + Inches(0.6), cw - 2*PAD, Inches(1.8), sz=Pt(13), clr=DARK, sp=Pt(12))

    bt = Inches(4.0); bh = Inches(2.5)
    card(s, MX, bt, SW - 2*MX, bh)
    txt(s, "Грант КОД-ИИ / ANIEMORE — подтверждённый результат",
        MX + PAD, bt + Inches(0.15), SW - 2*MX - 2*PAD, Inches(0.35),
        sz=Pt(17), clr=BLUE, bold=True)
    lines(s, [
        "Компания уже доказала свой высокий класс работы с государственными грантами",
        "В рамках КОД-ИИ разработана опенсорсная ИИ-библиотека для распознавания эмоций в речи человека — ANIEMORE",
        "Благодаря разработке ВУЗы и отделения РАН создали свои решения в Санкт-Петербурге, Новосибирске, Севастополе",
        "Корпорация МТС внедрила технологию в работу в своих сервисах",
        "Библиотека распространяется как опенсорс и уже стала стандартом для распознавания эмоций на русском языке",
    ], MX + PAD, bt + Inches(0.6), SW - 2*MX - 2*PAD, Inches(1.7), sz=Pt(13), clr=DARK, sp=Pt(10))

    txt(s, "Доказанный опыт работы с государственными грантами, контрактами и внедрениями на федеральном уровне",
        MX, Inches(6.85), SW - 2*MX, Inches(0.35), sz=Pt(13), clr=MID, al=PP_ALIGN.CENTER)


# =====================================================================
def main():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    for fn in [slide_01, slide_02, slide_03, slide_04, slide_05,
               slide_06, slide_07, slide_08, slide_09, slide_10,
               slide_11, slide_12]:
        fn(prs)
    prs.save(OUT)
    print("Сохранено:", OUT)


if __name__ == "__main__":
    main()
